import os
import random
import spacy
from transformers import pipeline
from pptx import Presentation
from docx import Document
import pdfplumber
from PIL import Image
import pytesseract

# --------------------------
# Load AI Models
# --------------------------
try:
    nlp = spacy.load("en_core_web_sm")
    text_generator = pipeline("text2text-generation", model="t5-small")
    print("✅ AI models loaded successfully!")
except Exception as e:
    print(f"❌ Error loading AI models: {e}")
    nlp = None
    text_generator = None

# --------------------------
# Extract text from file
# --------------------------
def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    text = ""
    try:
        if ext in [".txt", ".java"]:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == ".docx":
            doc = Document(filepath)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".pdf":
            with pdfplumber.open(filepath) as pdf:
                pages = []
                for page in pdf.pages:
                    t = page.extract_text()
                    if not t:
                        img = page.to_image(resolution=200).original
                        t = pytesseract.image_to_string(img)
                    pages.append(t)
                text = "\n".join(pages)
        elif ext == ".pptx":
            prs = Presentation(filepath)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            text = "\n".join(slides_text)
        elif ext in [".png", ".jpg", ".jpeg"]:
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img)
    except Exception as e:
        print(f"⚠️ Could not read text from {filepath}: {e}")
    return text.strip()

# --------------------------
# Extract entities (NER + few-shot)
# --------------------------
def extract_entities(text):
    entities = []
    if nlp:
        doc = nlp(text)
        entities = [ent.text for ent in doc.ents]
    if len(entities) < 5:
        words = [w for w in text.split() if len(w) > 3]
        entities += random.sample(words, min(10, len(words)))
    return list(dict.fromkeys(entities))

# --------------------------
# Generate meaningful quiz question
# --------------------------
def generate_quiz_question(entity, material):
    if not text_generator:
        return {"question": entity, "options":["A) Yes","B) No","C) Maybe","D) N/A"], "answer":"A"}
    
    prompt = f"""
Generate a meaningful multiple-choice question from this study material:
Material: {material}
Focus on this entity: {entity}
Provide one correct answer and three plausible distractors.
Format:
Question: ...
Answer: ...
Distractors: ...
"""
    try:
        gen = text_generator(prompt, max_new_tokens=300, do_sample=True)[0]["generated_text"]
        lines = [l.strip() for l in gen.splitlines() if l.strip()]
        question = next((l.split("Question:")[-1].strip() for l in lines if "Question:" in l), entity)
        answer = next((l.split("Answer:")[-1].strip() for l in lines if "Answer:" in l), entity)
        distractors_line = next((l.split("Distractors:")[-1].strip() for l in lines if "Distractors:" in l), "")
        distractors = [d.strip() for d in distractors_line.split(",") if d.strip()]
        while len(distractors) < 3:
            distractors.append(random.choice([w for w in material.split() if w.lower() != answer.lower()]))
        distractors = distractors[:3]

        # Shuffle options
        all_options = [answer] + distractors
        random.shuffle(all_options)
        letters = ["A","B","C","D"]
        options = [f"{letters[i]}) {all_options[i]}" for i in range(4)]
        correct_letter = letters[all_options.index(answer)]
        return {"question": question, "options": options, "answer": correct_letter}
    except:
        return {"question": entity, "options":["A) Yes","B) No","C) Maybe","D) N/A"], "answer":"A"}

# --------------------------
# Knowledge graph
# --------------------------
def build_knowledge_graph(entities):
    kg = {}
    for ent in entities:
        related = random.sample([e for e in entities if e!=ent], min(3, len(entities)-1))
        kg[ent] = related
    return kg

# --------------------------
# Main quiz loop
# --------------------------
def start_quiz():
    print("=== AI Study Buddy (MLV Mode) ===")
    choice = input("Do you want to (1) Paste text or (2) Upload file? ").strip()
    text_data = ""
    if choice=="1":
        text_data = input("Paste your text here:\n")
    elif choice=="2":
        filepath = input("Enter full path (txt, java, docx, pptx, pdf, jpg/png): ").strip()
        if not os.path.exists(filepath):
            print("⚠️ File not found!"); return
        text_data = extract_text_from_file(filepath)
        if not text_data:
            print("⚠️ Could not read any text from the file."); return
    else:
        print("Invalid choice!"); return

    print("Generating meaningful quiz and knowledge graph...")
    entities = extract_entities(text_data)
    knowledge_graph = build_knowledge_graph(entities)
    quiz_questions = [generate_quiz_question(ent, text_data) for ent in entities[:10]]

    print("Quiz Started!")
    for idx, q in enumerate(quiz_questions,1):
        print(f"\nQ{idx}: {q['question']}")
        for opt in q["options"]:
            print(opt)
        ans = input("Your answer (A/B/C/D or type 'exit' to quit): ").strip().upper()
        if ans=="EXIT": break
        if ans==q["answer"]:
            print("✅ Correct!")
        else:
            print(f"❌ Incorrect! Correct answer: {q['answer']})")

    print("\n🎉 Quiz Finished!\n")
    print("Knowledge Graph (Sample):")
    for k,v in knowledge_graph.items():
        print(f"{k} -> {v}")

if __name__=="__main__":
    start_quiz()
